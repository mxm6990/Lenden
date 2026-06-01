#!/usr/bin/env python3
"""Generate LenDen stakeholder PowerPoint deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# LenDen brand colors
BLACK = RGBColor(0x08, 0x0A, 0x09)
DARK = RGBColor(0x0F, 0x14, 0x12)
GREEN = RGBColor(0x1A, 0x5C, 0x38)
GREEN_DARK = RGBColor(0x0F, 0x3D, 0x24)
MINT = RGBColor(0x4A, 0xDE, 0x80)
MUTED = RGBColor(0x8A, 0x9A, 0x92)
WHITE = RGBColor(0xF8, 0xFA, 0xF9)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
CARD = RGBColor(0x1A, 0x22, 0x1E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.65)
MARGIN_R = Inches(0.65)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_glow(slide):
    add_rect(slide, Inches(9.5), Inches(-0.5), Inches(4), Inches(4), GREEN_DARK)


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=13, color=WHITE, bullet_color=MINT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.level = 0
        p.space_after = Pt(6)
        p.bullet = True
    return box


def add_eyebrow(slide, text):
    add_textbox(slide, MARGIN_L, Inches(0.45), CONTENT_W, Inches(0.3), text, size=9, bold=True, color=MINT)


def add_title(slide, text, top=Inches(0.75)):
    add_textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.9), text, size=28, bold=True, color=WHITE)


def add_stat_card(slide, left, top, width, height, value, label, source=None, highlight=False):
    fill = GREEN if highlight else CARD
    add_rect(slide, left, top, width, height, fill, MUTED if not highlight else MINT)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.55), value, size=24, bold=True, color=MINT)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.65), width - Inches(0.3), Inches(0.55), label, size=10, color=MUTED)
    if source:
        add_textbox(slide, left + Inches(0.15), top + height - Inches(0.35), width - Inches(0.3), Inches(0.25), source, size=8, color=MUTED)


def slide_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BLACK)
    add_glow(slide)

    add_rect(slide, MARGIN_L, Inches(0.55), Inches(0.55), Inches(0.55), GREEN, MINT)
    add_textbox(slide, MARGIN_L + Inches(0.12), Inches(0.58), Inches(0.4), Inches(0.5), "L", size=22, bold=True, color=MINT)

    add_textbox(slide, MARGIN_L + Inches(0.7), Inches(0.58), Inches(4), Inches(0.45), "LenDen  ·  লেন্দেন", size=20, bold=True, color=WHITE)

    add_rect(slide, MARGIN_L, Inches(1.25), Inches(3.2), Inches(0.32), CARD, MINT)
    add_textbox(slide, MARGIN_L + Inches(0.12), Inches(1.28), Inches(3), Inches(0.28), "STAKEHOLDER BRIEFING · MAY 2026", size=8, bold=True, color=MINT)

    add_textbox(slide, MARGIN_L, Inches(1.85), Inches(8), Inches(1.2), "Clarity in every investment.", size=40, bold=True, color=WHITE)
    add_textbox(slide, MARGIN_L, Inches(2.75), Inches(8), Inches(0.5), "প্রতিটি বিনিয়োগে স্পষ্টতা।", size=22, color=MUTED)

    add_textbox(
        slide, MARGIN_L, Inches(3.45), Inches(8.5), Inches(0.9),
        "A mobile-first retail investing experience for Bangladesh — built for clarity, trust, and regulatory alignment from day one.",
        size=14, color=MUTED,
    )

    y = Inches(4.55)
    for label, value in [
        ("Prepared for", "Bangladesh Bank · Government · Investors · IT Pioneers"),
        ("Status", "Closed-beta prototype · Paper trading only"),
        ("Contact", "mahathirmahbubimf@gmail.com"),
    ]:
        add_textbox(slide, MARGIN_L, y, Inches(1.4), Inches(0.25), label, size=10, bold=True, color=WHITE)
        add_textbox(slide, MARGIN_L + Inches(1.5), y, Inches(6.5), Inches(0.35), value, size=10, color=MUTED)
        y += Inches(0.42)


def slide_executive_summary(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "EXECUTIVE SUMMARY")
    add_title(slide, "Modern retail investing UX for a market that needs trust")

    items = [
        "LenDen is a bilingual (English / Bengali) mobile app prototype for DSE retail investing UX.",
        "Current version is a closed beta with mock trading — no real order routing, payments, or brokerage execution.",
        "Designed to partner with licensed brokers, Depository Participants (DPs), and regulated data providers.",
        "Targets the gap between 1.6M+ registered BO accounts and genuinely active, confident retail participation.",
    ]
    add_bullets(slide, MARGIN_L, Inches(1.75), Inches(7.2), Inches(2.8), items, size=12)

    cards = [
        ("Prototype", "Current stage — closed beta, 3–5 invited testers (internal target)"),
        ("Mock only", "No real money movement · No BSEC brokerage license claimed"),
        ("Mobile-first", "React + Capacitor iOS shell · Supabase backend"),
    ]
    x = Inches(8.1)
    y = Inches(1.75)
    for val, lbl in cards:
        add_stat_card(slide, x, y, Inches(4.5), Inches(1.15), val, lbl, highlight=(val == "Prototype"))
        y += Inches(1.3)

    add_rect(slide, MARGIN_L, Inches(5.85), CONTENT_W, Inches(0.85), CARD, AMBER)
    add_textbox(
        slide, MARGIN_L + Inches(0.15), Inches(5.95), CONTENT_W - Inches(0.3), Inches(0.7),
        "Transparency note: LenDen product metrics (users, AUM, revenue) are not included — the project is in prototype stage. All market statistics cite public sources.",
        size=10, color=WHITE,
    )


def slide_market_context(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "BANGLADESH CAPITAL MARKET · VERIFIED DATA")
    add_title(slide, "A large registered base — but shrinking active participation")

    stats = [
        ("16.39L", "Total BO accounts\n24 Dec 2025", "CDBL via TBS", True),
        ("12.05L", "BO accounts\nholding shares", "CDBL via TBS", False),
        ("3.67L", "BO zero-balance\n(+23,538 YoY)", "CDBL via TBS", False),
        ("~4–5L", "Estimated active traders\n(brokerage estimates)", "TBS Dec 2025", False),
    ]
    x = MARGIN_L
    for val, lbl, src, hl in stats:
        add_stat_card(slide, x, Inches(1.75), Inches(2.85), Inches(1.35), val, lbl, src, highlight=hl)
        x += Inches(3.05)

    add_textbox(slide, MARGIN_L, Inches(3.35), Inches(3), Inches(0.3), "BO account trend", size=12, bold=True, color=WHITE)
    bars = [("2016", "29.2L", 1.0), ("2024", "16.7L", 0.57), ("Dec 2025", "16.39L", 0.56)]
    y = Inches(3.75)
    for label, val, pct in bars:
        add_textbox(slide, MARGIN_L, y, Inches(1.2), Inches(0.25), label, size=10, color=MUTED)
        bar_w = Inches(5.5 * pct)
        add_rect(slide, MARGIN_L + Inches(1.3), y + Inches(0.04), bar_w, Inches(0.18), MINT if pct > 0.6 else AMBER)
        add_textbox(slide, MARGIN_L + Inches(7.0), y, Inches(1), Inches(0.25), val, size=10, bold=True, color=WHITE)
        y += Inches(0.38)

    add_textbox(
        slide, MARGIN_L, Inches(5.5), CONTENT_W, Inches(0.4),
        "1.26 million investors exited over 8 years (2016→2024). Source: Financial Express, citing DSE data.",
        size=9, color=MUTED,
    )


def slide_retail_decline(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "INVESTOR PROFILE · BSEC / CDBL DATA")
    add_title(slide, "Small retail investors are leaving faster than institutions return")

    distribution = [
        ("Portfolio < Tk 1L", "8.32L accounts", 0.85),
        ("Tk 1L – 10L", "2.85L", 0.35),
        ("Tk 10L – 50L", "80,608", 0.12),
        ("> Tk 1 crore", "13,316", 0.04),
    ]
    y = Inches(1.75)
    for label, val, pct in distribution:
        add_textbox(slide, MARGIN_L, y, Inches(1.8), Inches(0.25), label, size=10, color=MUTED)
        add_rect(slide, MARGIN_L + Inches(1.9), y + Inches(0.04), Inches(4.5 * pct), Inches(0.16), AMBER if pct > 0.7 else MINT)
        add_textbox(slide, MARGIN_L + Inches(6.6), y, Inches(1.2), Inches(0.25), val, size=10, bold=True, color=WHITE)
        y += Inches(0.38)

    add_textbox(slide, MARGIN_L, Inches(3.55), Inches(6), Inches(0.3), "Portfolio distribution as of Jun/Jul 2025. Sources: BSEC, TBS.", size=8, color=MUTED)

    cards = [
        ("37.9% → ~20%", "Retail shareholding in DSE-listed companies fell nearly half over 12 years", "Daily Star / DSE"),
        ("−9%", "Accounts below Tk 1L declined by 84,409 in one year (Jun 2024→Jun 2025)", "BSEC via TBS"),
        ("DSEX ~4,883", "Benchmark down 333 points in 2025; DS30 at ~1,882 (Dec 2025)", "TBS Dec 2025"),
    ]
    y = Inches(1.75)
    for val, lbl, src in cards:
        add_stat_card(slide, Inches(8.0), y, Inches(4.6), Inches(1.25), val, lbl, src, highlight=(val.startswith("37")))
        y += Inches(1.4)


def slide_digital(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "DIGITAL INFRASTRUCTURE · BTRC")
    add_title(slide, "Connectivity is national — capital markets UX is not")

    for i, (val, lbl, src) in enumerate([
        ("186M", "Mobile phone subscribers\nMarch 2026", "BTRC / AMTOB"),
        ("130M", "Total internet subscribers\nMarch 2026", "BTRC / AMTOB"),
        ("115M", "Mobile internet subscribers\nMarch 2026", "BTRC / AMTOB"),
    ]):
        add_stat_card(slide, MARGIN_L + Inches(i * 3.15), Inches(1.75), Inches(2.95), Inches(1.35), val, lbl, src, highlight=True)

    add_rect(slide, MARGIN_L, Inches(3.35), Inches(5.9), Inches(2.5), CARD, MUTED)
    add_textbox(slide, MARGIN_L + Inches(0.15), Inches(3.45), Inches(5.5), Inches(0.3), "The opportunity gap", size=13, bold=True, color=WHITE)
    add_bullets(slide, MARGIN_L + Inches(0.15), Inches(3.85), Inches(5.5), Inches(1.8), [
        "Hundreds of millions of Bangladeshis are already on smartphones and mobile data.",
        "Retail investing interfaces remain fragmented — broker portals and legacy terminals dominate.",
        "No dominant mobile-native, bilingual retail investing experience exists for the DSE today.",
    ], size=11)

    add_rect(slide, Inches(6.85), Inches(3.35), Inches(5.85), Inches(2.5), CARD, MINT)
    add_textbox(slide, Inches(7.0), Inches(3.45), Inches(5.5), Inches(0.3), "LenDen's design thesis", size=13, bold=True, color=WHITE)
    add_bullets(slide, Inches(7.0), Inches(3.85), Inches(5.5), Inches(1.8), [
        "Meet investors on the device they already use daily.",
        "Deliver clarity — portfolio, market, compliance — in Bengali and English.",
        "Build trust infrastructure (KYC, audit, disclosures) before scaling execution.",
    ], size=11)


def slide_product(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "PRODUCT")
    add_title(slide, "LenDen — clarity in every investment")

    add_textbox(
        slide, MARGIN_L, Inches(1.65), Inches(6.5), Inches(0.6),
        "LenDen (লেন্দেন) is a mobile-first prototype for exploring DSE retail investing UX. It is intentionally honest about its prototype status.",
        size=12, color=MUTED,
    )

    features = [
        "Home — Portfolio summary, buying power, DSE session status, watchlist",
        "Market — DSE index, search, securities list",
        "Portfolio — Holdings, performance chart, sector allocation",
        "Stock detail — Quotes, position, company context",
        "Trust Center — KYC flows, security, legal consents, support",
        "Mock buy flow — Order preview and simulated submit (paper trading)",
    ]
    add_bullets(slide, MARGIN_L, Inches(2.35), Inches(6.5), Inches(3.2), features, size=11)

    add_rect(slide, Inches(7.4), Inches(1.65), Inches(5.3), Inches(0.32), CARD, AMBER)
    add_textbox(slide, Inches(7.55), Inches(1.68), Inches(5), Inches(0.28), "PROTOTYPE BOUNDARIES", size=8, bold=True, color=AMBER)

    boundaries = [
        "No real order routing to DSE",
        "No real payment movement (bKash, bank transfers)",
        "No investment advice",
        "No licensed live market data in default mode",
        "Learn tab & quick actions — built, disabled for later",
    ]
    add_bullets(slide, Inches(7.4), Inches(2.1), Inches(5.3), Inches(2.5), boundaries, size=11)

    add_rect(slide, Inches(7.4), Inches(4.85), Inches(5.3), Inches(1.0), CARD, MUTED)
    add_textbox(slide, Inches(7.55), Inches(4.95), Inches(5), Inches(0.25), "DSE trading hours modeled", size=12, bold=True, color=WHITE)
    add_textbox(slide, Inches(7.55), Inches(5.25), Inches(5), Inches(0.5), "Sunday – Thursday · 10:00 – 14:30 (Asia/Dhaka). Mock orders respect market-open gate.", size=10, color=MUTED)


def slide_architecture(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "TECHNOLOGY · FOR IT PIONEERS")
    add_title(slide, "Modern, auditable, mobile-native architecture")

    blocks = [
        ("Client", "React 19 · TypeScript · Vite · Tailwind CSS v4 · Capacitor 7 (iOS)"),
        ("Backend", "Supabase (Auth, Postgres, Edge Functions) · Row-level security ready"),
        ("Market data modes", "Mock (default) · Experimental DSE (proxied, unofficial) · Licensed vendor (planned)"),
    ]
    y = Inches(1.75)
    for title, body in blocks:
        add_rect(slide, MARGIN_L, y, Inches(5.8), Inches(0.95), GREEN if title == "Client" else CARD, MINT if title == "Client" else MUTED)
        add_textbox(slide, MARGIN_L + Inches(0.15), y + Inches(0.08), Inches(5.4), Inches(0.25), title, size=12, bold=True, color=WHITE)
        add_textbox(slide, MARGIN_L + Inches(0.15), y + Inches(0.38), Inches(5.4), Inches(0.5), body, size=10, color=MUTED)
        y += Inches(1.1)

    add_textbox(slide, Inches(7.0), Inches(1.75), Inches(5.5), Inches(0.3), "Data flow (current prototype)", size=12, bold=True, color=WHITE)
    flow = ["Mobile App", "Supabase Auth", "Edge Functions", "Postgres + Cache"]
    x = Inches(7.0)
    for step in flow:
        add_rect(slide, x, Inches(2.15), Inches(1.35), Inches(0.55), CARD, MUTED)
        add_textbox(slide, x + Inches(0.05), Inches(2.25), Inches(1.25), Inches(0.4), step, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += Inches(1.45)

    add_rect(slide, Inches(7.0), Inches(3.0), Inches(5.5), Inches(2.4), CARD, MUTED)
    add_textbox(slide, Inches(7.15), Inches(3.1), Inches(5.2), Inches(0.25), "Security posture (planned)", size=12, bold=True, color=WHITE)
    add_bullets(slide, Inches(7.15), Inches(3.45), Inches(5.2), Inches(1.6), [
        "TLS in transit · encryption at rest for NID/TIN",
        "Append-only audit logs for sensitive actions",
        "Least-privilege admin · session/device management",
        "NID masked in UI; access logged",
    ], size=10)
    add_textbox(slide, Inches(7.15), Inches(5.05), Inches(5.2), Inches(0.3), "Note: Production encryption not yet implemented — prototype stores mock data.", size=8, color=MUTED)


def slide_regulatory(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "REGULATORY ALIGNMENT · BANGLADESH BANK & BSEC")
    add_title(slide, "Partner-first model — not an unlicensed broker")

    reqs = [
        "Licensed brokerage partner aligned with BSEC requirements",
        "Order routing to exchange or broker OMS",
        "Settlement & custody via CDBL",
        "BO account opening through a Depository Participant",
        "Licensed market data vendor (not community/unofficial feeds)",
        "AML/KYC provider integration and screening",
    ]
    add_textbox(slide, MARGIN_L, Inches(1.65), Inches(6), Inches(0.3), "What live execution would require", size=13, bold=True, color=WHITE)
    add_bullets(slide, MARGIN_L, Inches(2.0), Inches(6.2), Inches(3.5), reqs, size=11)

    add_rect(slide, Inches(7.2), Inches(1.65), Inches(5.5), Inches(1.35), GREEN, MINT)
    add_textbox(slide, Inches(7.35), Inches(1.75), Inches(5.2), Inches(0.25), "LenDen's intended role", size=12, bold=True, color=WHITE)
    add_textbox(
        slide, Inches(7.35), Inches(2.05), Inches(5.2), Inches(0.8),
        "An investing UX and education layer for Bangladeshi retail investors — not a standalone licensed broker on day one.",
        size=10, color=MUTED,
    )

    add_rect(slide, Inches(7.2), Inches(3.2), Inches(5.5), Inches(2.2), CARD, MUTED)
    add_textbox(slide, Inches(7.35), Inches(3.3), Inches(5.2), Inches(0.25), "Prototype compliance stance", size=12, bold=True, color=WHITE)
    add_bullets(slide, Inches(7.35), Inches(3.65), Inches(5.2), Inches(1.6), [
        'Clear "prototype / mock trading" labeling throughout app',
        "Risk disclosures and legal consents in Trust Center",
        "Market-hours enforcement on mock orders",
        "No misleading claims of live brokerage status",
    ], size=10)


def slide_investors(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "FOR INVESTORS")
    add_title(slide, "Addressing a trust deficit with product-led clarity")

    cols = [
        ("Problem", "Retail shareholding halved. BO accounts emptying. Market cap under pressure. Confidence gap is structural."),
        ("Solution", "A consumer-grade mobile experience that makes portfolio state, costs, and risks legible — rebuilding engagement through transparency."),
        ("Moat (planned)", "Bilingual UX · trust infrastructure · broker/DP integrations · licensed data · education layer."),
    ]
    x = MARGIN_L
    for title, body in cols:
        add_rect(slide, x, Inches(1.75), Inches(3.85), Inches(1.55), GREEN, MINT)
        add_textbox(slide, x + Inches(0.15), Inches(1.85), Inches(3.5), Inches(0.25), title, size=12, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.15), Inches(2.15), Inches(3.5), Inches(1.0), body, size=10, color=MUTED)
        x += Inches(4.05)

    add_rect(slide, MARGIN_L, Inches(3.55), Inches(5.9), Inches(2.3), CARD, MUTED)
    add_textbox(slide, MARGIN_L + Inches(0.15), Inches(3.65), Inches(5.5), Inches(0.25), "Business model options (not yet chosen)", size=12, bold=True, color=WHITE)
    add_bullets(slide, MARGIN_L + Inches(0.15), Inches(4.0), Inches(5.5), Inches(1.2), [
        "Brokerage revenue share with licensed partner",
        "Subscription for premium analytics / education",
        "B2B white-label UX for existing brokers",
    ], size=10)
    add_textbox(slide, MARGIN_L + Inches(0.15), Inches(5.15), Inches(5.5), Inches(0.35), "Not verified: Revenue model, unit economics, and fundraising terms are not finalized.", size=9, color=AMBER)

    add_rect(slide, Inches(6.85), Inches(3.55), Inches(5.85), Inches(2.3), CARD, MINT)
    add_textbox(slide, Inches(7.0), Inches(3.65), Inches(5.5), Inches(0.25), "Market sizing anchor (verified)", size=12, bold=True, color=WHITE)
    add_bullets(slide, Inches(7.0), Inches(4.0), Inches(5.5), Inches(1.4), [
        "~16.4L registered BO accounts (addressable registration base)",
        "~12L accounts with any shareholding",
        "DSE equity market cap ~Tk 3.1–3.5 trillion range (2025, varying by date/source)",
    ], size=10)


def slide_roadmap(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "ROADMAP")
    add_title(slide, "Phased path from prototype to regulated product")

    phases = [
        ("Phase 0 · Complete", "UX prototype & closed beta", "Mobile app, mock trading, Supabase auth, Trust Center, bilingual UI, admin concept dashboard"),
        ("Phase 1 · Next", "Expanded beta & broker conversations", "Define broker API contract · align BO onboarding with DP partner · licensed market data evaluation"),
        ("Phase 2", "Regulatory & integration", "KYC/AML provider · replace mock services · CDBL-aligned account flows · BSEC compliance review"),
        ("Phase 3", "Controlled launch", "Live execution via licensed partner · Learn tab · real payments · public rollout"),
    ]
    y = Inches(1.75)
    for phase, title, desc in phases:
        add_rect(slide, MARGIN_L, y, Inches(0.12), Inches(0.75), MINT if "Complete" in phase or "Next" in phase else MUTED)
        add_textbox(slide, MARGIN_L + Inches(0.25), y, Inches(2.2), Inches(0.25), phase, size=9, bold=True, color=MINT)
        add_textbox(slide, MARGIN_L + Inches(2.5), y, Inches(4), Inches(0.25), title, size=12, bold=True, color=WHITE)
        add_textbox(slide, MARGIN_L + Inches(2.5), y + Inches(0.3), Inches(9.5), Inches(0.4), desc, size=10, color=MUTED)
        y += Inches(0.95)

    add_textbox(slide, MARGIN_L, Inches(5.85), CONTENT_W, Inches(0.35), "Timeline is indicative — subject to regulatory approval and partnership timelines.", size=9, color=MUTED)


def slide_why_now(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "STRATEGIC TIMING")
    add_title(slide, "Why LenDen, why now")

    quadrants = [
        ("Market forces", [
            "Retail exodus creates demand for a better on-ramp when confidence returns",
            "Institutional/HNI participation rising — retail UX gap widens",
            "DSEX volatility underscores need for clear, honest investor communication",
        ]),
        ("Technology enablers", [
            "130M+ internet users — mobile-first is default",
            "Cloud auth/DB (Supabase) lowers time-to-beta",
            "Capacitor enables native iOS without dual codebases",
        ]),
        ("Policy alignment", [
            "Digital Bangladesh — financial inclusion via mobile",
            "Capital market deepening is a stated policy goal",
            "Regulator emphasis on BO account quality over raw count",
        ]),
        ("LenDen's ask", [
            "Guidance on regulatory pathway for UX-layer + broker partnership model",
            "Introductions to licensed brokers & DPs",
            "Investor/partner conversations for Phase 1–2 funding",
        ]),
    ]
    positions = [(MARGIN_L, Inches(1.75)), (Inches(6.85), Inches(1.75)), (MARGIN_L, Inches(4.15)), (Inches(6.85), Inches(4.15))]
    for (left, top), (title, bullets) in zip(positions, quadrants):
        highlight = title == "LenDen's ask"
        add_rect(slide, left, top, Inches(5.85), Inches(2.1), GREEN if highlight else CARD, MINT if highlight else MUTED)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.12), Inches(5.5), Inches(0.25), title, size=12, bold=True, color=WHITE)
        add_bullets(slide, left + Inches(0.15), top + Inches(0.45), Inches(5.5), Inches(1.5), bullets, size=10)


def slide_sources(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    add_glow(slide)
    add_eyebrow(slide, "APPENDIX · DATA INTEGRITY")
    add_title(slide, "Verified sources & explicit unknowns")

    rows = [
        ("BO accounts (total)", "16.39 lakh (24 Dec 2025)", "CDBL via TBS"),
        ("BO with shares", "12.05 lakh", "CDBL via TBS"),
        ("BO zero balance", "3.67 lakh", "CDBL via TBS"),
        ("Accounts < Tk 1L", "8.32 lakh (Jun 2025)", "BSEC via TBS"),
        ("Accounts > Tk 1Cr", "13,316 (Jun 2025)", "BSEC via TBS"),
        ("Retail shareholding", "37.9% → ~20% (12 yrs)", "Daily Star / DSE"),
        ("BO decline 2016→2024", "29.2L → 16.7L", "Financial Express / DSE"),
        ("Mobile subscribers", "186.06M (Mar 2026)", "BTRC / AMTOB"),
        ("Internet subscribers", "129.62M (Mar 2026)", "BTRC / AMTOB"),
        ("DSEX (Dec 2025)", "~4,883 (−333 YTD)", "TBS"),
    ]

    y = Inches(1.65)
    add_textbox(slide, MARGIN_L, y, Inches(3.5), Inches(0.22), "Metric", size=8, bold=True, color=MINT)
    add_textbox(slide, MARGIN_L + Inches(3.6), y, Inches(3.5), Inches(0.22), "Value", size=8, bold=True, color=MINT)
    add_textbox(slide, MARGIN_L + Inches(7.2), y, Inches(4.5), Inches(0.22), "Source", size=8, bold=True, color=MINT)
    y += Inches(0.28)

    for metric, value, source in rows:
        add_textbox(slide, MARGIN_L, y, Inches(3.5), Inches(0.22), metric, size=8, bold=True, color=WHITE)
        add_textbox(slide, MARGIN_L + Inches(3.6), y, Inches(3.5), Inches(0.22), value, size=8, color=MUTED)
        add_textbox(slide, MARGIN_L + Inches(7.2), y, Inches(4.5), Inches(0.22), source, size=8, color=MUTED)
        y += Inches(0.24)

    add_rect(slide, MARGIN_L, Inches(4.35), CONTENT_W, Inches(1.1), CARD, AMBER)
    add_textbox(
        slide, MARGIN_L + Inches(0.15), Inches(4.45), CONTENT_W - Inches(0.3), Inches(0.95),
        "Not included (unknown / not applicable): LenDen user count, transaction volume, AUM, revenue, funding raised, team size, broker partnership status, BSEC/Bangladesh Bank approval status, bKash integration metrics, or competitor market share.",
        size=9, color=WHITE,
    )


def slide_thank_you(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BLACK)
    add_glow(slide)

    add_rect(slide, Inches(6.2), Inches(1.5), Inches(0.9), Inches(0.9), GREEN, MINT)
    add_textbox(slide, Inches(6.38), Inches(1.62), Inches(0.6), Inches(0.7), "L", size=32, bold=True, color=MINT, align=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN_L, Inches(2.65), CONTENT_W, Inches(0.8), "Thank you", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN_L, Inches(3.45), CONTENT_W, Inches(0.45), "প্রতিটি বিনিয়োগে স্পষ্টতা।", size=18, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(
        slide, MARGIN_L, Inches(3.95), CONTENT_W, Inches(0.5),
        "Let's build trustworthy retail investing for Bangladesh — together.",
        size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )

    meta_y = Inches(4.85)
    for label, value in [("Email", "mahathirmahbubimf@gmail.com"), ("Product", "LenDen · Closed Beta Prototype"), ("Deck version", "May 2026 · Stakeholder v1")]:
        add_textbox(slide, MARGIN_L, meta_y, Inches(4.2), Inches(0.22), label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(4.5), meta_y, Inches(4.2), Inches(0.22), value, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        meta_y += Inches(0.35)


def build():
    prs = new_prs()
    slide_title(prs)
    slide_executive_summary(prs)
    slide_market_context(prs)
    slide_retail_decline(prs)
    slide_digital(prs)
    slide_product(prs)
    slide_architecture(prs)
    slide_regulatory(prs)
    slide_investors(prs)
    slide_roadmap(prs)
    slide_why_now(prs)
    slide_sources(prs)
    slide_thank_you(prs)

    out = "/Users/mahathirmahbub/Desktop/LenDen/presentation/LenDen-Stakeholder-Deck.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
